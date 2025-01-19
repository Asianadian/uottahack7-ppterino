from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from llm import LLM
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.dml.color import RGBColor
import random
import requests
from io import BytesIO

load_dotenv()
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
llm = LLM(OPENAI_API_KEY)

app = Flask(__name__)
CORS(app)

ppt_code = """
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
"""

def generateSlide(template_num, title, content_arr):
    global ppt_code
    if template_num == 1:
        ppt_code += f"\nslide = prs.slides.add_slide(prs.slide_layouts[1])"
        ppt_code += f"\ntitle = slide.shapes.title"
        ppt_code += f"\ncontent = slide.placeholders[1]"

        ppt_code += f"\ntitle.text = '{title}'"
        
        result_string = 'points = [\n    ' + ',\n    '.join(f'"{point}"' for point in content_arr) + '\n]'
        ppt_code += f"\n{result_string}"

        ppt_code += f"""
for point in points:
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = point
"""

    elif template_num == 2:
        ppt_code += f"\nslide = prs.slides.add_slide(prs.slide_layouts[2])"
        ppt_code += f"\ntitle = slide.shapes.title"
        ppt_code += f"\ntitle.text = '{title}'"
    elif template_num == 3:
        ppt_code += f"\nslide = prs.slides.add_slide(prs.slide_layouts[3])"
        ppt_code += f"\ntitle = slide.shapes.title"
        ppt_code += f"\ntitle.text = '{title}'"

        # ppt_code += f"\left_placeholder = slide.placeholders[1]"
        # ppt_code += f"\nleft_placeholder.text = 'Here is some text for the left column.'"

        # ppt_code += f"\right_placeholder = slide.placeholders[2]"
        # ppt_code += f"\right_placeholder.text = 'Here is some text for the right column.'"  
    elif template_num == 8:
        ppt_code += f"\nslide = prs.slides.add_slide(prs.slide_layouts[8])"
        ppt_code += f"\ntitle = slide.shapes.title"
        ppt_code += f"\ntitle.text = '{title}'"


@app.route("/api/pptx", methods=['POST'])
def pptx():
    global ppt_code

    try:
        # getting request data from frontend
        data = request.form
        print(data)
        prompt = data.get('prompt')
        script = data.get('script')
        imageUrl = data.get('imageUrl')

        files = request.files
        transcription = files.get('file')
        try:
            save_path = os.path.join("uploads", transcription.filename)
            os.makedirs("uploads", exist_ok=True)  # Create 'uploads' directory if it doesn't exist
            transcription.save(save_path)
        except:
            ...
        print('transcribing')
        transcription = llm.transcribe(save_path)

        image = files.get('image')
        try:
            save_path = os.path.join("uploads", image.filename)
            os.makedirs("uploads", exist_ok=True)  # Create 'uploads' directory if it doesn't exist
            image.save(save_path)
        except:
            ...
        print('describing')
        image_desc = llm.describe_image(save_path)

        prompt += ' ' + transcription + ' ' + image_desc

        ppt = None

        print('done')

        # generating summarized points and storing in parsed_slide_text
        parsed_slide_text = []
        cur_title = ""
        cur_points = []
        title = ""

        slide_text = llm.gen_slides_from_script(script)

        for line in slide_text.splitlines():
            if line.startswith('#### '):
                title = line[5:].strip()
            elif line.startswith('### '):
                # Save the previous section if exists
                if cur_title:
                    parsed_slide_text.append((cur_title, cur_points.copy()))

                # Start a new title section
                cur_title = line[4:].strip()
                cur_points = []
            elif line.startswith('## '):
                cur_points.append(line[3:].strip())

        # generating code
        ppt_code += f"\ntitle.text = '{title}'"

        for title, points in parsed_slide_text:
            number = random.choice([1, 2, 3, 8])
            generateSlide(number, title, points)
        
        ppt_code += f"\nprs.save('res.pptx')"

        while True:
            try:
                print(ppt_code)
                exec(ppt_code)
                
                ppt = Presentation('./res.pptx') 
                break
            except:
                break

        # editing ppt slide backgrounds and titles
        slide_width = ppt.slide_width
        slide_height = ppt.slide_height

        possible_backgrounds = ["test.png", "image.jpeg", "imagetest.png"]

        # INTRODUCE LOGIC HERE, IF WE PROVIDE AN IMAGE WE FEED THAT AND GENERATE VARIATIONS
        # IF NO IMAGE, WE RUN THE GENERATION FIRST

        if imageUrl != "":
            # variation_url = llm.generate_variations(imageUrl)
            # resp = requests.get(variation_url)
            # image_stream = BytesIO(resp.content)
            # possible_backgrounds.append(image_stream)
            pass
        else:
            # image_url = llm.create_ppt_bg(prompt)
            # possible_backgrounds.append(image_url)
            pass

        for i in range(2):
            # image_url = "test.png"
            # variation_url = llm.generate_variations(image_url)

            # resp = requests.get(variation_url)
            # image_stream = BytesIO(resp.content)
            # possible_backgrounds.append(image_stream)
            pass
        print(ppt)
        for slide in ppt.slides:
            print('generating')
            image_url = llm.create_ppt_bg(prompt)
            print(image_url)
            resp = requests.get(image_url)
            image_stream = BytesIO(resp.content)
            print('appending')
            pic = slide.shapes.add_picture(image_stream, 0, 0, width=slide_width, height=slide_height)
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

        ppt.save("res.pptx")
        res = send_file('./res.pptx')
        return res
        
    except Exception as e:
        return {"error": f"Failed to create events: {str(e)}"}

app.run(debug=True)